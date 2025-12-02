# Copyright (c) 2024, DjaoDjin inc.
# see LICENSE.

import csv, datetime, io, logging, os, re

from django.conf import settings
from django.db import connection
from django.db.models import Q, F
from django.http import HttpResponse
from django.shortcuts import get_object_or_404
from django.views.generic import ListView, TemplateView
from openpyxl import Workbook
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.shapes.autoshape import Shape
from pptx.shapes.graphfrm import GraphicFrame
from rest_framework.generics import ListAPIView
from rest_framework.request import Request
from survey.api.matrix import (AccessiblesAccountsMixin, BenchmarkMixin,
    EngagedAccountsMixin)
from survey.compat import force_str, six
from survey.filters import DateRangeFilter
from survey.helpers import datetime_or_now, get_extra, extra_as_internal
from survey.mixins import TimersMixin
from survey.models import (Answer, Campaign, Choice, Portfolio, Sample, Unit,
    UnitEquivalences)
from survey.settings import DB_PATH_SEP
from survey.utils import (get_accessible_accounts, get_engaged_accounts,
    get_question_model)

from .base import CSVDownloadRenderer
from .content import PracticesSpreadsheetView
from .. import humanize
from ..api.campaigns import CampaignContentMixin
from ..api.serializers import EngagementSerializer, LongFormatSerializer
from ..api.portfolios import (BenchmarkMixin as BenchmarkPSPMixin,
    CompletionRateMixin,
    DashboardAggregateMixin, EngagementStatsMixin,
    PortfolioAccessibleSamplesMixin, PortfolioEngagementMixin)
from ..compat import gettext_lazy as _
from ..helpers import as_valid_sheet_title
from ..mixins import (AccountMixin, CampaignMixin,
    AccountsNominativeQuerysetMixin)
from ..models import ScorecardCache
from ..scores.base import get_top_normalized_score
from ..utils import get_alliances, get_practice_serializer

LOGGER = logging.getLogger(__name__)

STAGE_BY_BUCKET = {
    '1-40': _("Adopting/ Initiating"),
    '41-60': _("Growing/ Progressing"),
    '61-80': _("Leading"),
    '81-100': _("Pioneering"),
}

def get_bucket(normalized_score):
    bucket = "1-40"
    if normalized_score:
        if normalized_score > 40.5:
            bucket = "41-60"
        if normalized_score > 60.5:
            bucket = "61-80"
        if normalized_score > 80.5:
            bucket = "81-100"
    return bucket


class CSVDownloadView(AccountMixin, ListView):

    content_type = 'text/csv'
    basename = 'download'
    headings = []
    filter_backends = []

    @staticmethod
    def encode(text):
        if six.PY2:
            return force_str(text).encode('utf-8')
        return force_str(text)

    @staticmethod
    def decorate_queryset(queryset):
        return queryset

    def filter_queryset(self, queryset):
        """
        Recreating a GenericAPIView.filter_queryset functionality here
        """
        # creating a DRF-compatible request object
        request = Request(self.request)
        for backend in list(self.filter_backends):
            queryset = backend().filter_queryset(request, queryset, self)
        return queryset

    def get(self, *args, **kwargs): #pylint: disable=unused-argument
        if six.PY2:
            content = io.BytesIO()
        else:
            content = io.StringIO()
        #pylint:disable=attribute-defined-outside-init
        self.csv_writer = csv.writer(content)
        self.csv_writer.writerow([self.encode(head)
            for head in self.get_headings()])
        qs = self.decorate_queryset(self.filter_queryset(self.get_queryset()))
        for record in qs:
            self.writerecord(record)
        content.seek(0)
        resp = HttpResponse(content, content_type=self.content_type)
        resp['Content-Disposition'] = 'attachment; filename="{}"'.format(
            self.get_filename())
        return resp

    def get_headings(self):
        return self.headings

    def get_queryset(self):
        # Note: this should take the same arguments as for
        # Searchable and SortableListMixin in "extra_views"
        raise NotImplementedError

    def get_filename(self):
        return datetime_or_now().strftime(self.basename + '-%Y%m%d.csv')

    def queryrow_to_columns(self, record):
        raise NotImplementedError

    def writerecord(self, record):
        table = self.queryrow_to_columns(record)
        if table:
            if isinstance(table[0], list):
                for row in table:
                    self.csv_writer.writerow(row)
            else:
                self.csv_writer.writerow(table)


class FullReportPPTXView(CampaignMixin, AccountMixin, TemplateView):
    """
    Download full report as a .pptx presentation
    """
    content_type = \
     'application/vnd.openxmlformats-officedocument.presentationml.presentation'
    basename = 'report'
    is_percentage = True
    title = 'Full Report'

    def get_data(self, title=None):
        #pylint:disable=unused-argument
        return []

    def get_filename(self):
        return datetime_or_now().strftime(
            self.account.slug + '-' + self.basename + '-%Y%m%d.pptx')

    def get_questions_by_key(self, prefix=None, initial=None):
        #pylint:disable=unused-argument
        return initial if isinstance(initial, dict) else {}

    def get_template_names(self):
        candidates = [
            os.path.join('app', 'reporting', '%s.pptx' % self.campaign),
            os.path.join('app', 'reporting', '%s.pptx' % self.basename)]
        return candidates + super(FullReportPPTXView, self).get_template_names()

    def get(self, request, *args, **kwargs):
        #pylint: disable=unused-argument,too-many-locals,too-many-nested-blocks
        context = {
            'title': self.title,
            'accounts': ', '.join([self.account.printable_name] + [
                alliance.printable_name for alliance
                in get_alliances(self.account)]),
            'unit': "%" if self.is_percentage else "nb"
        }

        # Prepares the result file
        content = io.BytesIO()
        candidate = None
        for candidate_template in self.get_template_names():
            for template_dir in settings.TEMPLATES_DIRS:
                full_path = os.path.join(template_dir, candidate_template)
                if os.path.exists(full_path):
                    candidate = full_path
                    break
            if candidate:
                break
        LOGGER.debug("use template '%s'", candidate)
        if candidate:
            with open(candidate, 'rb') as reporting_file:
                prs = Presentation(reporting_file)
            for slide in prs.slides:
                LOGGER.debug("slide=%s", slide)
                for shape in slide.shapes:
                    LOGGER.debug("\tshape=%s", shape)
                    if isinstance(shape, Shape):
                        LOGGER.debug("\t- text=%s", shape.text)
                        for key, val in context.items():
                            key_text = '{{%s}}' % key
                            if key_text in shape.text:
                                for para in shape.text_frame.paragraphs:
                                    LOGGER.debug("\t- replace %s in %s",
                                        key_text, para)
                                    para.text = para.text.replace(key_text, val)
                    elif isinstance(shape, GraphicFrame):
                        # We found the chart's container
                        try:
                            chart = shape.chart
                            if chart.has_title:
                                data = list(self.get_data(chart.chart_title))
                            else:
                                data = list(self.get_data())
                            LOGGER.debug("use data %s", data)
                            chart_data = CategoryChartData()
                            # Series might not have exactly the same labels.
                            # Thus we need to gather all defined labels first.
                            labels = set([])
                            for serie in data:
                                for point in serie.get('values'):
                                    label = point[0]
                                    if isinstance(label, datetime.datetime):
                                        label = label.date()
                                    labels |= {label}
                            labels = sorted(labels)
                            LOGGER.debug("\tlabels=%s", labels)
                            for label in labels:
                                chart_data.add_category(label)
                            # Make sure we have a value for each label before
                            # adding serie.
                            LOGGER.debug("\tseries=%s", list(data))
                            for serie in data:
                                dataset = {point[0]:point[1]
                                    for point in serie.get('values')}
                                values = []
                                for label in labels:
                                    val = dataset.get(label)
                                    values += [val if val else 0]
                                LOGGER.debug("\tadd serie '%s' with values %s",
                                    serie.get('title'), values)
                                chart_data.add_series(
                                    serie.get('title'), values)

                            chart.replace_data(chart_data)
                        except ValueError:
                            pass
            prs.save(content)
        content.seek(0)

        resp = HttpResponse(content, content_type=self.content_type)
        resp['Content-Disposition'] = 'attachment; filename="{}"'.format(
            self.get_filename())

        return resp


class BenchmarkPPTXView(BenchmarkPSPMixin, FullReportPPTXView):
    """
    Download benchmarks a .pptx presentation
    """
    basename = 'benchmarks'
    serializer_class = get_practice_serializer()

    @property
    def question(self):
        if not hasattr(self, '_question'):
            #pylint:disable=attribute-defined-outside-init
            self._question = get_object_or_404(
                get_question_model(), path=self.db_path)
        return self._question

    def get_template_names(self):
        if self.question.default_unit.system == Unit.SYSTEM_DATETIME:
            return [os.path.join('app', 'reporting', 'column.pptx')]
        return [os.path.join('app', 'reporting', 'doughnut.pptx')]

    @property
    def title(self):
        #pylint:disable=attribute-defined-outside-init
        if not hasattr(self, '_title'):
            self._title = self.question.title
        return self._title

    def get_title(self):
        # To label benchmark data set (engaged suppliers)
        return self.account.printable_name

    def get_query_param(self, key, default_value=None):
        # XXX Workaround until overridden method in djaodjin-survey supports
        # Django HTTPRequest.
        try:
            return self.request.query_params.get(key, default_value)
        except AttributeError:
            pass
        return self.request.GET.get(key, default_value)

    def get_decorated_questions(self, prefix=None):
        return list(six.itervalues(self.get_questions_by_key(
            prefix=prefix if prefix else DB_PATH_SEP)))

    def get_data(self, title=None):
        questions = self.get_decorated_questions(self.db_path)
        results = questions[0].get('benchmarks', [])
        if self.question.default_unit.system != Unit.SYSTEM_DATETIME:
            return reversed(results)
        return results


class ReportingDashboardPPTXView(DashboardAggregateMixin, FullReportPPTXView):
    """
    Download reporting dashboard as an .pptx spreadsheet.
    """
    basename = 'doughnut'

    def get_template_names(self):
        return [os.path.join('app', 'reporting', '%s.pptx' % self.basename)]

    def get_filename(self):
        return datetime_or_now().strftime(
            self.account.slug + '-' + self.basename + '-%Y%m%d.pptx')

    def get_data(self, title=None):
        data = self.get_response_data(self.request, **self.kwargs)
        return data['results']

    @property
    def title(self):
        #pylint:disable=attribute-defined-outside-init
        if not hasattr(self, '_title'):
            data = self.get_response_data(self.request, **self.kwargs)
            self._title = data['title']
        return self._title


class CompletionRatePPTXView(CompletionRateMixin, ReportingDashboardPPTXView):
    """
    Download completion rate as a .pptx presentation
    """
    basename = 'completion-rate'


class EngagementStatsPPTXView(EngagementStatsMixin, ReportingDashboardPPTXView):
    """
    Download engagement statistics as a .pptx presentation
    """

    def get_response_data(self, request, *args, **kwargs):
        resp = super(EngagementStatsPPTXView, self).get_response_data(
            request, *args, **kwargs)
        # We have to reverse the results to keep charts consistent with
        # HTML version.
        resp.update({'results': reversed(resp.get('results'))})
        return resp


class TemplateXLSXView(TimersMixin, AccountMixin, ListView):

    content_type = \
        'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'

    basename = 'download'
    filter_backends = []
    title = ""

    def __init__(self, **kwargs):
        super(TemplateXLSXView, self).__init__(**kwargs)
        self.wsheet = None

    def decorate_queryset(self, queryset):
        return queryset

    def filter_queryset(self, queryset):
        """
        Given a queryset, filter it with whichever filter backend is in use.
        """
        for backend in list(self.filter_backends):
            queryset = backend().filter_queryset(self.request, queryset, self)
        return queryset

    def get_descr(self):
        return ""

    def get_filename(self):
        return datetime_or_now().strftime(
            self.account.slug + '-' + self.basename + '-%Y%m%d.xlsx')

    def get_headings(self):
        return []

    @staticmethod
    def get_indent_question(depth=0):
        return "  " * depth

    @staticmethod
    def get_indent_heading(depth=0):
        return "  " * depth

    def get(self, request, *args, **kwargs):
        #pylint: disable=unused-argument
        self._start_time()
        wbook = Workbook()

        # Populate the Total sheet
        self.wsheet = wbook.active
        self.wsheet.title = as_valid_sheet_title(self.title)

        queryset = self.filter_queryset(self.get_queryset())
        self.decorate_queryset(queryset)

        descr = self.get_descr()
        if descr:
            self.wsheet.append([descr])
        self.wsheet.append(self.get_headings())
        self.write_queryset(queryset)

        # Prepares the result file
        content = io.BytesIO()
        wbook.save(content)
        content.seek(0)

        resp = HttpResponse(content, content_type=self.content_type)
        resp['Content-Disposition'] = 'attachment; filename="{}"'.format(
            self.get_filename())
        self._report_queries("http response created")
        return resp

    def write_queryset(self, queryset):
        for record in queryset:
            self.writerecord(record)

    def writerecord(self, record):
        table = self.queryrow_to_columns(record)
        if table:
            if isinstance(table[0], list):
                for row in table:
                    self.wsheet.append(row)
            else:
                self.wsheet.append(table)

    def queryrow_to_columns(self, record):
        raise NotImplementedError


class PortfolioAccessiblesXLSXView(PortfolioAccessibleSamplesMixin,
                                   TemplateXLSXView):
    """
    Download scores of all reporting entities as an Excel spreadsheet.
    """
    basename = 'accessibles'
    title = "Accessibles responses"

    def get_headings(self):
        # External PowerBI workflows depend on both, "Supplier No." and
        # "Supplier Name" being capitalized as such.
        labels = [str(_("Supplier No.")), str(_("Supplier Name")),
            str(_("Tags"))]
        for label in list(reversed(self.labels)):
            labels += ["%s Score" % label]
        for label in list(reversed(self.labels)):
            labels += ["%s Status" % label]
        for label in list(reversed(self.labels)):
            labels += ["%s Completed at" % label]
        for label in list(reversed(self.labels)):
            labels += ["%s Score Range" % label]
        return labels

    def queryrow_to_columns(self, record):
        reporting_statuses = dict(humanize.REPORTING_STATUSES)
        supplier_key = get_extra(record, 'supplier_key')
        row = [str(supplier_key), record.printable_name,
            ", ".join(get_extra(record, 'tags', []))]

        # Write period-over-period normalized scores
        for val in reversed(record.values):
            state = val.state
            if state in (humanize.REPORTING_COMPLETED,
                         humanize.REPORTING_VERIFIED):
                normalized_score = get_top_normalized_score(
                    val, segments_candidates=self.segments_candidates)
                row += [
                    normalized_score if normalized_score is not None else ""]
            else:
                row += [""]

        # Write period-over-period status
        for val in reversed(record.values):
            state = val.state
            row += [reporting_statuses[state]]

        # Write period-over-period completed at
        for val in reversed(record.values):
            completed_at = val.created_at if val.pk else None
            row += [completed_at.strftime('%Y-%m-%d') if completed_at else ""]

        # Write period-over-period score buckets
        for val in reversed(record.values):
            state = val.state
            if state in (humanize.REPORTING_COMPLETED,
                         humanize.REPORTING_VERIFIED):
                normalized_score = get_top_normalized_score(
                    val, segments_candidates=self.segments_candidates)
                bucket = get_bucket(normalized_score)
                row += [bucket]
            else:
                row += [""]

        return row


class PortfolioAccessiblesLongCSVView(PortfolioAccessibleSamplesMixin,
                                      TemplateXLSXView):
    """
    Downloads year-over-year scores in long format
    """
    base_headings = ['Supplier No.', 'Supplier Name',
        'TSP Score', 'Score Range', 'Score Stage']
    title = "Responses"

    def get_headings(self):
        headings = self.base_headings
        if self.period == 'monthly':
            headings += ['Month']
        else:
            headings += ['Year']
        return headings

    def get_filename(self):
        basename = "%s-suppliers" % self.account
        return datetime_or_now().strftime(basename + '-%Y%m%d.xlsx')

    def queryrow_to_columns(self, record):
        collapsed_reporting_statuses = {
            humanize.REPORTING_NO_RESPONSE: _("Invited"),
            humanize.REPORTING_NO_DATA: _("No Data"),
            humanize.REPORTING_NO_PROFILE: _("N/A"),
            humanize.REPORTING_RESPONDED: _("Not Shared"),
        }
        reporting_statuses = dict(humanize.REPORTING_STATUSES)
        table = []
        for val in reversed(record.values):
            bucket = None
            stage = None
            # Write 'Supplier No.' and 'Supplier Name'
            supplier_key = get_extra(record, 'supplier_key')
            row = [str(supplier_key), record.printable_name]
            state = val.state
            # Write 'TSP Score'
            if state in (humanize.REPORTING_COMPLETED,
                         humanize.REPORTING_VERIFIED):
                normalized_score = get_top_normalized_score(
                    val, segments_candidates=self.segments_candidates)
                bucket = get_bucket(normalized_score)
                stage = STAGE_BY_BUCKET.get(bucket)
                row += [
                    normalized_score if normalized_score is not None else ""]
            else:
                status = collapsed_reporting_statuses.get(
                    state, reporting_statuses[state])
                bucket = status
                stage = status
                row += [str(status)]

            # Write 'Score Range'
            row += [str(bucket)]
            # Write 'Score Stage'
            row += [str(stage)]
            # Write Period ('Year' or 'Month')
            if self.period == 'monthly':
                row += [val.created_at.strftime("%b %Y")]
            else:
                row += [val.created_at.year]

            table += [row]

        return table


class PortfolioEngagementXLSXView(PortfolioEngagementMixin, TemplateXLSXView):
    """
    Download scores of all reporting entities as an Excel spreadsheet.
    """
    basename = 'engagement'
    title = "Engagement"

    def get_headings(self):
        return [
            'SupplierID',
            'Profile name',
            'Tags',
            'Status',
            'Last update',
            'Last reminder',
            'Added']

    def queryrow_to_columns(self, record):
        reporting_statuses = dict(EngagementSerializer.REPORTING_STATUSES)
        supplier_key = get_extra(record, 'supplier_key')
        row = [
            supplier_key,
            record.printable_name,
            ", ".join(get_extra(record, 'tags', [])),
            reporting_statuses[record.reporting_status],
            record.last_activity_at.date() if record.last_activity_at else "",
            "",
            record.requested_at.date() if record.requested_at else ""]
        return row


class LongFormatCSVView(AccountsNominativeQuerysetMixin, CSVDownloadView):
    """
    Download raw data in format suitable for OLAP Software
    """
    headings = ['Created at', 'SupplierID', 'Profile name',
        'Measured', 'Unit', 'Question title']
    filter_backends = [DateRangeFilter]

    def __init__(self, **kwargs):
        super(LongFormatCSVView, self).__init__(**kwargs)
        self._choices = {}

    def get_supplier_key(self, account):
        #pylint:disable=attribute-defined-outside-init
        if not hasattr(self, '_portfolio_extras'):
            extras_queryset = Portfolio.objects.filter(
                grantee=self.account,
                account_id__in=[val.pk for val in self.requested_accounts])
            self._portfolio_extras = {
                (val.grantee_id, val.account_id): extra_as_internal(val)
                for val in extras_queryset}
        return self._portfolio_extras.get(
            (self.account.pk, account.pk), {}).get('supplier_key', "")


    def get_queryset(self):
        if not self.requested_accounts:
            return Answer.objects.none()
        requested_accounts = self.requested_accounts.filter(
            grant_key__isnull=True)
        queryset = Answer.objects.filter(
            Q(unit_id=F('question__default_unit_id')) |
        Q(unit_id=F('question__default_unit__source_equivalences__target_id')),
            sample__in=Sample.objects.get_latest_frozen_by_accounts(
                campaign=self.campaign, accounts=requested_accounts,
                start_at=self.start_at, ends_at=self.ends_at, tags=[]
            )).distinct().select_related('sample__account', 'question', 'unit')
        return queryset

    def queryrow_to_columns(self, record):
        measured = record.measured
        if record.unit.system == Unit.SYSTEM_ENUMERATED:
            if not hasattr(self, '_choices'):
                self._choices = {}
            choices = self._choices.get(record.unit.slug)
            if not choices:
                choices = {choice.pk: choice.text
                    for choice in Choice.objects.filter(
                        unit=record.unit)}
                self._choices.update({record.unit.slug: choices})
            measured = choices.get(record.measured)
        elif record.unit.system not in Unit.NUMERICAL_SYSTEMS:
            measured = self.encode(
                Choice.objects.get(pk=measured).text)
        supplier_key = self.get_supplier_key(record.sample.account)
        row = [
            record.created_at.date(),
            supplier_key,
            record.sample.account.printable_name,
            measured,
            record.unit.slug,
            self.encode(record.question.title),
            record.question.path
        ]
        return row


class AnswersDownloadMixin(BenchmarkMixin, CampaignContentMixin, TimersMixin):
    """
    Download answers and scores for a set of accounts
    """
    #pylint:disable=too-many-instance-attributes
    basename = 'answers'
    show_comments = True
    add_expanded_styles = False
    paginator = None

#    ordering = ('full_name',)
    ordering = ('account_id',)

    path_idx = 0
    account_id_idx = 1
    measured_idx = 2
    unit_idx = 3
    default_unit_idx = 4
    unit_title_idx = 5
    measured_text_idx = 6

    def __init__(self, *args):
        super(AnswersDownloadMixin, self).__init__( *args)
        self.comments_unit = Unit.objects.get(slug='freetext')
        self.points_unit = Unit.objects.get(slug='points')
        self.target_by_unit = Unit.objects.get(slug='ends-at')
        self.errors = []

    @property
    def search_terms(self):
        if not hasattr(self, '_search_terms'):
            self._search_terms = self.get_query_param('q', None)
        return self._search_terms

    @property
    def show_planned(self):
        if not hasattr(self, '_show_planned'):
            self._show_planned =self.get_query_param('planned', False)
            if not isinstance(self._show_planned, bool):
                try:
                    self._show_planned = bool(int(self._show_planned))
                except ValueError:
                    self._show_planned = bool(
                        self._show_planned.lower() in ['true'])
        return self._show_planned

    @property
    def show_scores(self):
        if not hasattr(self, '_show_scores'):
            self._show_scores =self.get_query_param('scores', True)
            if not isinstance(self._show_scores, bool):
                try:
                    self._show_scores = bool(int(self._show_scores))
                except ValueError:
                    self._show_scores = bool(
                        self._show_scores.lower() not in ['false'])
        return self._show_scores

    @property
    def verified_campaign(self):
        if not hasattr(self, '_verified_campaign'):
            #pylint:disable=attribute-defined-outside-init
            self._verified_campaign = self.campaign
            look = re.match(r'(\S+)-verified$', self._verified_campaign.slug)
            if look:
                self._verified_campaign = get_object_or_404(
                    Campaign.objects.all(), slug=look.group(1))
        return self._verified_campaign

    @property
    def latest_assessments(self):
        if not hasattr(self, '_latest_assessments'):
            #pylint:disable=attribute-defined-outside-init
            if self.accounts_with_engagement:
                self._latest_assessments = \
                    Sample.objects.get_latest_frozen_by_accounts(
                        campaign=self.verified_campaign,
                        segment_prefix=self.db_path, ends_at=self.ends_at,
                    # create a list to prevent RawQuerySet if-condition later on
                        grantees=[self.account],
                        accounts=[account.pk for account
                            in self.accounts_with_engagement],
                        tags=[])
            else:
                self._latest_assessments = Sample.objects.none()
        return self._latest_assessments

    @property
    def latest_improvements(self):
        if not hasattr(self, '_improvements'):
            #pylint:disable=attribute-defined-outside-init
            if self.accounts_with_engagement:
                self._improvements = \
                    Sample.objects.get_latest_frozen_by_accounts(
                        campaign=self.verified_campaign,
                        segment_prefix=self.db_path, ends_at=self.ends_at,
                        grantees=[self.account],
                # create a list to prevent RawQuerySet if-condition later on
                        accounts=[account.pk for account
                            in self.accounts_with_engagement],
                        tags=['is_planned'])
            else:
                self._improvements = Sample.objects.none()
        return self._improvements

    @property
    def accounts_with_engagement(self):
        #pylint:disable=attribute-defined-outside-init
        if not hasattr(self, '_accounts_with_engagement'):
            self._accounts_with_engagement = list(self.get_accounts())
            extras = {
                (val.grantee_id, val.account_id): extra_as_internal(val)
                for val in Portfolio.objects.filter(
                    grantee=self.account, # XXX forces a single grantee
                    account_id__in=[
                        val.pk for val in self.accounts_with_engagement])}
            for reporting in self._accounts_with_engagement:
                # Merge portfolio extra field into account extra field.
                extra = extras.get((self.account.pk, reporting.pk))
                if extra:
                    reporting.extra = extra_as_internal(reporting)
                    reporting.extra.update(extra)
        return self._accounts_with_engagement


    def merge_records(self, left_records, right_records):
        """
        Merge two lists of tuples ordered by 'account_id' (i.e. at index 1)

        The returned list is ordered 'account_id', and contains all records
        in the initial `left_records` and `right_records` lists.
        """
        results = []
        left_iterator = iter(left_records)
        right_iterator = iter(right_records)
        try:
            left = next(left_iterator)
        except StopIteration:
            left = None
        try:
            right = next(right_iterator)
        except StopIteration:
            right = None
        try:
            while left and right:
                left_account_id = left[self.account_id_idx]
                right_account_id = right[self.account_id_idx]
                if left_account_id < right_account_id:
                    results += [left]
                    left = None
                    left = next(left_iterator)
                elif left_account_id > right_account_id:
                    results += [right]
                    right = None
                    right = next(right_iterator)
                else:
                    # equals
                    results += [left]
                    results += [right]
                    left = None
                    right = None
                    try:
                        left = next(left_iterator)
                    except StopIteration:
                        pass
                    try:
                        right = next(right_iterator)
                    except StopIteration:
                        pass
        except StopIteration:
            pass
        try:
            while left:
                results += [left]
                left = None
                left = next(left_iterator)
        except StopIteration:
            pass
        try:
            while right:
                results += [right]
                right = None
                right = next(right_iterator)
        except StopIteration:
            pass

        return results


    def get_answers_by_paths(self, latest_samples, prefix=None):
        """
        Returns a dictionnary of lists of answers tuple (path, account_id,
        measured, unit_id, default_unit_id, text) by question slug

             {
               question_slug: [
                 (path, account_id, measured, unit_id, default_unit_id, title,
                  text),
                 ...
               ]
             }
        """
        answers_by_paths = {}
        try:
            # We might have a `RawQuerySet` so we can't blindly use `.exists()`
            unused_first_sample = latest_samples[0]
        except IndexError:
            return answers_by_paths

        question_clause = ""
        if prefix:
            question_clause = ("WHERE survey_question.path LIKE '%s%%'"
                % prefix)
        reporting_answers_sql = """
WITH samples AS (
    %(latest_assessments)s
),
answers AS (
SELECT
    survey_question.path,
    samples.account_id,
    survey_answer.measured,
    survey_answer.unit_id,
    survey_question.default_unit_id,
    survey_unit.title
FROM survey_answer
INNER JOIN survey_question
  ON survey_answer.question_id = survey_question.id
INNER JOIN samples
  ON survey_answer.sample_id = samples.id
INNER JOIN survey_unit
  ON survey_answer.unit_id = survey_unit.id
%(question_clause)s
)
SELECT
  answers.path,
  answers.account_id,
  answers.measured,
  answers.unit_id,
  answers.default_unit_id,
  answers.title,
  survey_choice.text
FROM answers
LEFT OUTER JOIN survey_choice
  ON survey_choice.unit_id = answers.unit_id AND
     survey_choice.id = answers.measured
ORDER BY answers.path, answers.account_id
""" % {
            'latest_assessments': latest_samples.query.sql,
            'question_clause': question_clause
        }
        if self.campaign != self.verified_campaign:
            # When we are downloading the answers for a verification
            # campaign, we run the same query as previously, except we
            # add an indirection through `djaopsp_verifiedsample` to find
            # the account the verification answers apply to (otherwise
            # we would set the verifier account as a column label).
            reporting_answers_sql = """
WITH latest_samples AS (
    %(latest_assessments)s
),
samples AS (
SELECT survey_sample.id,
 survey_sample.slug,
 survey_sample.created_at,
 survey_sample.campaign_id,
 latest_samples.account_id AS account_id,
 survey_sample.is_frozen,
 survey_sample.extra,
 survey_sample.updated_at
FROM survey_sample
INNER JOIN djaopsp_verifiedsample
ON djaopsp_verifiedsample.verifier_notes_id = survey_sample.id
INNER JOIN latest_samples
ON djaopsp_verifiedsample.sample_id = latest_samples.id
WHERE djaopsp_verifiedsample.verified_status >= 2
),
answers AS (
SELECT
    survey_question.path,
    samples.account_id,
    survey_answer.measured,
    survey_answer.unit_id,
    survey_question.default_unit_id,
    survey_unit.title
FROM survey_answer
INNER JOIN survey_question
  ON survey_answer.question_id = survey_question.id
INNER JOIN samples
  ON survey_answer.sample_id = samples.id
INNER JOIN survey_unit
  ON survey_answer.unit_id = survey_unit.id
%(question_clause)s
)
SELECT
  answers.path,
  answers.account_id,
  answers.measured,
  answers.unit_id,
  answers.default_unit_id,
  answers.title,
  survey_choice.text
FROM answers
LEFT OUTER JOIN survey_choice
  ON survey_choice.unit_id = answers.unit_id AND
     survey_choice.id = answers.measured
ORDER BY answers.path, answers.account_id
""" % {
            'latest_assessments': latest_samples.query.sql,
            'question_clause': question_clause
        }

        with connection.cursor() as cursor:
            cursor.execute(reporting_answers_sql, params=None)
            prev_path = None
            chunk = []
            for row in cursor:
                # The SQL quesry is ordered by `path` so we can build
                # the final result by chunks, path by path.
                path = row[0]
                if prev_path and path != prev_path:
                    # flush
                    prev_slug = prev_path.split(DB_PATH_SEP)[-1]
                    if prev_slug in answers_by_paths:
                        # Ideally the query would be sorted by slug
                        # i.e. `REGEXP_REPLACE(answers.path, '^/(\S+/)*', '')`
                        # but this function is not supported easily on sqlite3.
                        answers_by_paths.update({
                            prev_slug: self.merge_records(
                                answers_by_paths[prev_slug], chunk)})
                    else:
                        answers_by_paths.update({prev_slug: chunk})
                    chunk = []
                    #self._report_queries(
                    #    "merged answers by slug %s" % prev_slug)
                chunk += [row]
                prev_path = path
            if chunk:
                # flush last remaining chunk
                prev_slug = prev_path.split(DB_PATH_SEP)[-1]
                if prev_slug in answers_by_paths:
                    answers_by_paths.update({
                        prev_slug: self.merge_records(
                            answers_by_paths[prev_slug], chunk)})
                else:
                    answers_by_paths.update({prev_slug: chunk})
                self._report_queries("merged answers by slug %s" % prev_slug)

        return answers_by_paths


    def get_queryset(self):
        """
        Returns a set of practices and headings decorated
        with answers per account
        """
        # Since we use `CampaignContentMixin` instead of `CampaignDecorateMixin`
        # directly, the set of questions will be current ones.
        questions = super(AnswersDownloadMixin, self).get_queryset()
        # Since we use `CampaignContentMixin` instead of `CampaignDecorateMixin`
        # directly, the set of questions will be current ones for the campaign.
        #
        # `questions` is a list of dictionnary formatted as:
        #     {
        #       "path": string,
        #       "rank": integer,
        #       "required": boolean,
        #       "default_unit": {
        #           "slug": string,
        #           "title": string,
        #           "system": integer
        #       },
        #       "title": string,
        #       "picture": string,
        #       "extra": dict,
        #     }

        # We build a tabularized list of answers per account for each question.
        #
        # `by_paths` is a dictionnary formatted as:
        #     {
        #       question_slug: [
        #         (path, account_id, measured, unit_id, default_unit_id, text),
        #         ...
        #       ]
        #     }
        #
        # We overide `get_queryset` instead of `get_decorated_questions`
        # and don't pass a `prefix` to `get_answers_by_paths` such that
        # answers are colapsed by 'slug'.
        if self.show_planned:
            by_paths = self.get_answers_by_paths(self.latest_improvements)
            self._report_queries(descr="collected planned answers")
        else:
            by_paths = self.get_answers_by_paths(self.latest_assessments)
            self._report_queries(descr="collected questions answers")

        for entry in questions:
            slug = entry.get('slug')
            if not slug:
                # We are running code here before `CampaignDecorateMixin`
                # has had a chance to add the slug.
                slug = entry.get('path').split(DB_PATH_SEP)[-1]
            by_accounts = by_paths.get(slug)
            if by_accounts:
                entry.update({'accounts': by_accounts})
            elif self.show_scores:
                # Headings are only added in `CampaignDecorateMixin` after
                # `get_decorated_questions` is called.
                tags = []
                extra = entry.get('extra')
                if extra:
                    tags = extra.get('tags')
                if tags and 'scorecard' in tags:
                    scores = ScorecardCache.objects.filter(
                        sample__in=self.latest_assessments,
                        path=entry['path']).order_by(
                        'sample__account_id').values_list(
                        'pk', 'sample__account_id', 'normalized_score')
                    # by using `self.points_unit.pk` for both 'unit_id' and
                    # 'default_unit_id', `add_datapoint` will set the 'measured'
                    # field.
                    entry.update({'accounts': [(
                        val[0],                    # path_idx = 0
                        val[1],                    # account_id_idx = 1
                        val[2],                    # measured_idx = 2
                        self.points_unit.pk,       # unit_idx = 3
                        self.points_unit.pk,       # default_unit_idx = 4
                        self.points_unit.title,    # unit_title_idx = 5
                        "")                        # measured_text_idx = 6
                        for val in scores]})
        self._report_queries(descr="collected section scores")

        return questions


class AnswersPivotableView(AnswersDownloadMixin, ListAPIView):

    basename = 'answers'
    headings = ['Created at', 'SupplierID', 'Profile name',
        'Measured', 'Unit', 'Question title', 'Question RefNum']

    def get_filename(self, ext='.csv'):
        return datetime_or_now().strftime(self.basename + '-%Y%m%d' + ext)

    def get_queryset(self):
        #pylint:disable=too-many-locals
        queryset = super(AnswersPivotableView, self).get_queryset()

        results = []
        last_activity_at_by_accounts = {}
        for sample in self.latest_assessments:
            last_activity_at_by_accounts.update({
                sample.account_id: sample.created_at})

        by_account_ids = {account.pk: account
            for account in self.accounts_with_engagement}

        for entry in queryset:
            question_title = entry.get('title')
            ref_num = entry.get('ref_num')
            by_accounts = entry.get('accounts', [])
            if by_accounts:
                for col in by_accounts:
                    account_id = col[self.account_id_idx]
                    account = by_account_ids[account_id]
                    supplier_key = get_extra(account, 'supplier_key')
                    last_activity_at = last_activity_at_by_accounts.get(
                        account_id)
                    measured = (col[self.measured_text_idx]
                        if col[self.measured_text_idx]
                        else col[self.measured_idx])
                    unit_title = col[self.unit_title_idx]
                    results += [{
                        'created_at': last_activity_at,
                        'supplier_key': supplier_key,
                        'printable_name': account.printable_name,
                        'measured': measured,
                        'unit': unit_title,
                        'title': question_title,
                        'ref_num': ref_num
                    }]

        return results


    def get(self, request, *args, **kwargs):
        self._start_time()
        return super(AnswersPivotableView, self).get(request, *args, **kwargs)


class AccessiblesAnswersPivotableCSVView(AccessiblesAccountsMixin,
                                         AnswersPivotableView):
    """
    Download answers and scores in format suitable for OLAP Software
    """
    basename = 'track-answers-pivotable'
    renderer_classes = [CSVDownloadRenderer]
    serializer_class = LongFormatSerializer

    def get_accounts(self):
        """
        Returns account accessibles by a profile in a specific date range.
        """
        # Redefines `AccessiblesAccountsMixin.get_accounts` such that
        # we use `self.verified_campaign` instead of `self.campaign`
        # and set `aggregate_set` to `False` instead of `True`.
        # Furthermore, we use the same [start_at, ends_at[ range
        # to filter accounts and samples.
        results = get_accessible_accounts([self.account],
            campaign=self.verified_campaign,
            start_at=self.start_at, ends_at=self.ends_at,
            aggregate_set=False).order_by('pk')
        return results


class EngagedAnswersPivotableCSVView(EngagedAccountsMixin,
                                     AnswersPivotableView):
    """
    Download answers and scores in format suitable for OLAP Software
    """
    basename = 'engage-answers-pivotable'
    renderer_classes = [CSVDownloadRenderer]
    serializer_class = LongFormatSerializer

    def get_accounts(self):
        """
        Returns account accessibles by a profile in a specific date range.
        """
        # Redefines `AccessiblesAccountsMixin.get_accounts` such that
        # we use `self.verified_campaign` instead of `self.campaign`
        # and set `aggregate_set` to `False` instead of `True`.
        # Furthermore, we use the same [start_at, ends_at[ range
        # to filter accounts and samples.
        return get_engaged_accounts([self.account],
            campaign=self.verified_campaign,
            start_at=self.start_at, ends_at=self.ends_at,
            aggregate_set=False,
            search_terms=self.search_terms).order_by('pk')


class TabularizedAnswersXLSXView(AnswersDownloadMixin,
                                 PracticesSpreadsheetView):
    """
    Download a spreadsheet of answers/comments with questions as rows
    and accessible accounts as columns.
    """

    @staticmethod
    def _get_title(element):
        title = element.get('title', "")
        ref_num = element.get('ref_num')
        if ref_num:
            title = "%s %s" % (ref_num, title)
        default_unit = element.get('default_unit', {})
        if default_unit and default_unit.get('system') in Unit.METRIC_SYSTEMS:
            title += " (in %s)" % default_unit.get('title')
        return title

    def add_error(self, msg):
        self.errors += [msg]

    def add_datapoint(self, account, row):
        """
        `account` is a dictionnary that contains answers with a field name
        by type/unit.

            {
              "account_id": integer,
              "measured": type,
              "comments": string,
              "score": integer
            }

        `row` is a (path, account_id, measured, unit_id, default_unit_id, text)
        tuple.
        """
        # XXX check if we are overriding a previous value since we are using
        # question slug instead of path.
        account_id = row[self.account_id_idx]
        measured = row[self.measured_idx]
        unit_id = row[self.unit_idx]
        default_unit_id = row[self.default_unit_idx]
        text = row[self.measured_text_idx]
        if unit_id == default_unit_id:
            # if we have already resolve `measured` to text value
            # (ex: because the unit is an enum), let's use that.
            if unit_id == self.target_by_unit.pk:
                try:
                    text = int(text)
                except ValueError:
                    pass
            if unit_id == self.comments_unit.pk:
                account.update({'measured': 'answered' if text else ""})
                account['comments'] += str(text) if text else ""
            else:
                account.update({'measured': text if text else measured})
        elif unit_id == self.points_unit.pk:
            account.update({'score': measured})
        elif unit_id == self.comments_unit.pk:
            account['comments'] += str(text) if text else ""
        else:
            unit = Unit.objects.get(pk=unit_id)
            if unit.system in (Unit.SYSTEM_STANDARD, Unit.SYSTEM_IMPERIAL):
                default_unit = Unit.objects.get(pk=default_unit_id)
                try:
                    equiv = UnitEquivalences.objects.get(
                        source=unit, target=default_unit)
                    account.update({
                        'measured': equiv.as_target_unit(measured)})
                except UnitEquivalences.DoesNotExist:
                    account.update({'measured': ""})
                    self.add_error("cannot convert '%s' to '%s' for %s:%s" % (
                        unit, default_unit, account_id, row[self.path_idx]))


    def as_account(self, key):
        """
        Fills a cell for column `key` with generated content since we do not
        have an item for it.
        """
        return {
            'account_id': key.pk, 'measured': "", 'comments': "", 'score': ""}

    @staticmethod
    def before_account(datapoint, account):
        """
        returns True if left < right
        """
        account_id = datapoint[1]
        return account_id < account.get('account_id')

    @staticmethod
    def after_account(datapoint, account):
        """
        returns True if left > right
        """
        account_id = datapoint[1]
        return account_id > account.get('account_id')

    @staticmethod
    def equals(aggregate, account):
        return aggregate.get('account_id') == account.get('account_id')


    def tabularize(self, datapoints, keys):
        """
        `datapoints` is a list of (path, account_id, measured, unit_id,
        default_unit_id, title, text) where path is constant, sorted
        by account_id.
        `keys` is a list of accounts sorted using the same account_id criteria.

        This function returns a sorted list of (account_id, {values}) where
        for each account in `keys`, there is a matching account in the results,
        inserting empty values where necessary.
        """
        results = []
        keys_iterator = iter(keys)
        datapoints_iterator = iter(datapoints)
        try:
            datapoint = next(datapoints_iterator)
            # path = datapoint[0]
        except StopIteration:
            datapoint = None
        prev_key = None
        try:
            key = self.as_account(next(keys_iterator))
        except StopIteration:
            key = None
        try:
            while datapoint and key:
                if not results or not self.equals(results[-1], key):
                    results += [key]
                if self.before_account(datapoint, key):
                    # This condition should never hold as datapoints' keys
                    # are a subset of keys, correct?
                    raise RuntimeError("Impossibility for "\
                      " before_account(datapoint=%s, key=%s)."\
                      " Maybe datapoints and/or keys are not sorted properly."
                      % (datapoint, key))
                if self.after_account(datapoint, key):
                    prev_key = key
                    key = None
                    key = self.as_account(next(keys_iterator))
                    if prev_key and prev_key == key:
                        raise RuntimeError(
                            "We have a duplicate key for %s" % str(key))
                else:
                    # equals
                    account = results[-1]
                    self.add_datapoint(account, datapoint)
                    try:
                        datapoint = next(datapoints_iterator)
                    except StopIteration:
                        datapoint = None
        except StopIteration:
            pass
        # We should have gone through all the datapoints since their
        # keys are a subset of `keys`.
        assert datapoint is None
        try:
            while key:
                if not results or not self.equals(results[-1], key):
                    results += [key]
                prev_key = key
                key = self.as_account(next(keys_iterator))
                if prev_key and prev_key == key:
                    raise RuntimeError(
                        "We have a duplicate key for %s" % str(key))
        except StopIteration:
            pass
        assert len(results) == len(keys)
        return results


    def format_row(self, entry, key=None):
        row = [self._get_title(entry)]
        by_accounts = self.tabularize(entry.get('accounts', []),
            self.accounts_with_engagement)
        if by_accounts:
            for account in by_accounts:
                # XXX Showing scores as % needs to be done when applying styles.
                row += [account.get(key, "")]
        else:
            for account in self.accounts_with_engagement:
                row += [""]

        if entry.get('default_unit'):
            slug = entry.get('slug')
            row += [slug]

        return row


    def get(self, request, *args, **kwargs):
        #pylint: disable=unused-argument,too-many-locals

        # We need to run `get_queryset` before `get_headings` so we know
        # how many columns to display for implementation rate.
        self._start_time()
        self.errors = []
        queryset = self.get_queryset()
        self._report_queries("built list of questions")

        key = 'measured'
        if self.show_planned:
            title = "Planned"
        else:
            title = "Answers"
        self.write_sheet(title=title, key=key, queryset=queryset)
        if self.show_comments:
            self.write_sheet(title="Comments", key='comments',
                queryset=queryset)

        if self.errors:
            LOGGER.info('\n'.join(self.errors))

        resp = HttpResponse(self.flush_writer(), content_type=self.content_type)
        resp['Content-Disposition'] = \
            'attachment; filename="{}"'.format(self.get_filename())
        return resp


    def write_headers(self):
        """
        Write table headers in the worksheet
        """
        ends_at = self.ends_at.date()
        source = "Source: %s" % self.request.build_absolute_uri(location='/')
        if not self.start_at:
            self.writerow(["Organization/profiles engaged to %s (%s)" % (
                ends_at.isoformat(), source)])
        else:
            start_at = self.start_at.date()
            self.writerow(["Organization/profiles engaged from %s to %s (%s)" %
                (start_at.isoformat(), ends_at.isoformat(), source)])
        headings = [force_str(_("Profile name"))] + [
            reporting.printable_name
            for reporting in self.accounts_with_engagement]
        self.writerow(headings)

        # Adds SupplierID and Tags headings
        key_headings = [force_str(_("SupplierID"))]
        tags_headings = [force_str(_("Tags"))]
        for reporting in self.accounts_with_engagement:
            key_headings += [get_extra(reporting, 'supplier_key')]
            tags = get_extra(reporting, 'tags')
            if tags:
                tags_headings += [','.join(tags)]
            else:
                tags_headings += [""]
        self.writerow(key_headings)
        self.writerow(tags_headings)

        # creates row with last completed date
        last_activity_at_by_accounts = {}
        for sample in self.latest_assessments:
            last_activity_at_by_accounts.update({
                sample.account_id: sample.created_at})
        headings = [force_str(_("Last completed at"))]
        for reporting in self.accounts_with_engagement:
            account_id = reporting.pk
            last_activity_at = last_activity_at_by_accounts.get(account_id)
            if last_activity_at:
                headings += [datetime.date(
                    last_activity_at.year,
                    last_activity_at.month,
                    last_activity_at.day)]
            else:
                headings += [""]
        self.writerow(headings)


class AccessiblesAnswersXLSXView(AccessiblesAccountsMixin,
                                 TabularizedAnswersXLSXView):
    """
    Download a spreadsheet of answers/comments with questions as rows
    and accessible accounts as columns.
    """
    basename = 'track-answers'

    def get_accounts(self):
        """
        Returns account accessibles by a profile in a specific date range.
        """
        # Redefines `AccessiblesAccountsMixin.get_accounts` such that
        # we use `self.verified_campaign` instead of `self.campaign`
        # and set `aggregate_set` to `False` instead of `True`.
        # Furthermore, we use the same [start_at, ends_at[ range
        # to filter accounts and samples.
        results = get_accessible_accounts([self.account],
            campaign=self.verified_campaign,
            start_at=self.start_at, ends_at=self.ends_at,
            aggregate_set=False).order_by('pk')
        return results


class EngagedAnswersXLSXView(EngagedAccountsMixin,
                             TabularizedAnswersXLSXView):
    """
    Download a spreadsheet of answers/comments with questions as rows
    and engaged accounts as columns.
    """
    basename = 'engage-answers'

    def get_accounts(self):
        """
        Returns account accessibles by a profile in a specific date range.
        """
        # Redefines `AccessiblesAccountsMixin.get_accounts` such that
        # we use `self.verified_campaign` instead of `self.campaign`
        # and set `aggregate_set` to `False` instead of `True`.
        # Furthermore, we use the same [start_at, ends_at[ range
        # to filter accounts and samples.
        results = get_engaged_accounts([self.account],
            campaign=self.verified_campaign,
            start_at=self.start_at, ends_at=self.ends_at,
            aggregate_set=False,
            search_terms=self.search_terms).order_by('pk')
        return results

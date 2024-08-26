# Copyright (c) 2024, DjaoDjin inc.
# see LICENSE.
from __future__ import unicode_literals

import datetime, json, logging
from collections import OrderedDict
from io import BytesIO

from deployutils.apps.django.templatetags.deployutils_prefixtags import (
    site_url)
from deployutils.helpers import update_context_urls
from django import forms
from django.contrib.auth import get_user_model
from django.core.files.storage import get_storage_class
from django.db import models, transaction
from django.http import Http404, HttpResponse, HttpResponseRedirect
from django.shortcuts import get_object_or_404
from django.views.generic import ListView
from django.views.generic.base import TemplateView
from django.views.generic.edit import FormMixin
from django.template.defaultfilters import slugify
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from survey.helpers import datetime_or_now, get_extra
from survey.models import (Campaign, Choice, EditableFilter,
    PortfolioDoubleOptIn, Sample)
from survey.settings import DB_PATH_SEP, URL_PATH_SEP
from survey.utils import get_account_model, get_question_model

from .downloads import PracticesSpreadsheetView
from ..scores import get_score_calculator
from ..api.samples import AssessmentContentMixin
from ..compat import reverse, six, gettext_lazy as _
from ..mixins import AccountMixin, SectionReportMixin
from ..utils import (get_latest_active_assessments,
    get_latest_completed_assessment)

LOGGER = logging.getLogger(__name__)


class AssessPracticesView(SectionReportMixin, TemplateView):
    """
    Profile assessment page
    """
    template_name = 'app/assess/index.html'
    breadcrumb_url = 'assess_practices'
    breadcrumb_url_index = 'assess_index'

    def get_reverse_kwargs(self):
        """
        List of kwargs taken from the url that needs to be passed through
        to ``reverse``.
        """
        return [self.account_url_kwarg, self.sample_url_kwarg,
            self.path_url_kwarg]

    def get_template_names(self):
        candidates = ['app/assess/%s.html' % self.sample.campaign]
        candidates += super(AssessPracticesView, self).get_template_names()
        return candidates

    def get_context_data(self, **kwargs):
        context = super(AssessPracticesView, self).get_context_data(**kwargs)
        context.update({
            'prefix': self.full_path,
            'nb_answers': self.nb_answers,
            'nb_questions': self.nb_questions,
            'units': {}
        })
        for unit_slug in ('verifiability', 'supporting-document',
                          'completeness'):
            context['units'].update({
                unit_slug.replace('-', '_'): Choice.objects.filter(
                    unit__slug=unit_slug).order_by('rank')})
        if not self.sample.is_frozen:
            context.update({
                'nb_required_answers': self.nb_required_answers,
                'nb_required_questions': self.nb_required_questions,
            })
        update_context_urls(context, {
            'pages_index': reverse('pages_index'),
            'track_metrics_index': reverse(
                'track_metrics_index', args=(self.account,)),
            'api_profiles': site_url("/api/accounts/users"),
            'api_assessment_sample': reverse('survey_api_sample',
                args=(self.account, self.sample)),
            'api_asset_upload_complete': self.request.build_absolute_uri(
                reverse('pages_api_upload_asset', args=(self.account,))),
            'api_aggregate_metric_base': reverse(
                'survey_api_aggregate_metric_base', args=(self.account,)),
            # These URLs can't be accessed by profiles the sample was shared
            # with. They must use ``sample.account``.
            'assess_base': reverse('assess_practices',
                args=(self.account, self.sample, '-'))[:-2],
            # XXX should download PDF with actions guidance.
            'print': reverse('assess_download_index', args=(
                self.account, self.sample)),
        })
        if self.path:
            url_path = self.path.lstrip(URL_PATH_SEP)
            api_path = self.full_path.lstrip(URL_PATH_SEP)
            update_context_urls(context, {
                'download': reverse('assess_download_segment', args=(
                    self.account, self.sample, url_path)),
                'api_content': reverse('api_sample_content', args=(
                    self.account, self.sample, api_path)),
                'api_account_benchmark': reverse(
                    'survey_api_sample_benchmarks',
                    args=(self.account, self.sample, api_path)),

            })
        else:
            update_context_urls(context, {
                'download': reverse('assess_download_index', args=(
                    self.account, self.sample)),
                'api_content': reverse('api_sample_content_index', args=(
                    self.account, self.sample)),
                'api_account_benchmark': reverse(
                    'survey_api_sample_benchmarks_index',
                    args=(self.account, self.sample)),
            })
        # Upload supporting documents
        storage_class = get_storage_class()
        if 's3boto' in storage_class.__name__.lower():
            update_context_urls(context, {
                'api_asset_upload_start': site_url(
                    "api/auth/tokens/realms/%s" % self.account),
            })
        else:
            update_context_urls(context, {
            'api_asset_upload_start': self.request.build_absolute_uri(
                reverse('pages_api_upload_asset', args=(self.account,))),
            })
        return context


class AssessRedirectForm(forms.Form):

    campaign = forms.CharField()


class AssessRedirectView(AccountMixin, FormMixin, TemplateView):
    """
    Shows a list of assessment requested or in progress decorated
    with grantees whenever available.
    """
    campaign_url_kwarg = 'campaign'
    path_url_kwarg = 'path'

    # XXX This code was copy/pasted from ScorecardRedirectView and mixed
    # with SourcingAssessRedirectView.
    template_name = 'app/assess/redirects.html'
    form_class = AssessRedirectForm

    def create_sample(self, campaign):
        account_model = get_account_model()
        with transaction.atomic():
            #pylint:disable=unused-variable
            if isinstance(self.account, account_model):
                account = self.account
            else:
                account, unused = account_model.objects.get_or_create(
                    slug=str(self.account))
            # XXX Whenever Sample.campaign_id is null, the survey APIs
            # will not behave properly.
            sample, created = Sample.objects.get_or_create(
                account=account, campaign=campaign, is_frozen=False,
                extra__isnull=True)
        return sample

    def form_valid(self, form):
        try:
            campaign = self.campaign_candidates.get(
                slug=form.cleaned_data['campaign'])
        except Campaign.DoesNotExist:
            raise Http404('No candidate campaign matches %(campaign)s.' % {
                'campaign': form.cleaned_data['campaign']})
        sample = self.create_sample(campaign)
        kwargs = {
            self.account_url_kwarg: self.account,
            'sample': sample
        }
        return HttpResponseRedirect(self.get_redirect_url(**kwargs))

    def get_redirect_url(self, *args, **kwargs):
        reverse_kwargs = {
            self.account_url_kwarg: kwargs.get(self.account_url_kwarg),
            'sample': kwargs.get('sample')
        }
        path = kwargs.get(self.path_url_kwarg)
        if path:
            reverse_kwargs.update({self.path_url_kwarg: path})
            return reverse('assess_practices', kwargs=reverse_kwargs)
        return reverse('assess_index', kwargs=reverse_kwargs)

    def get_context_data(self, **kwargs):
        context = super(AssessRedirectView, self).get_context_data(**kwargs)

        at_time = datetime_or_now()
        campaign_slug = self.kwargs.get(self.campaign_url_kwarg)
        campaign_filtered = (get_object_or_404(Campaign, slug=campaign_slug)
            if campaign_slug else None)
        path_filtered = self.kwargs.get(self.path_url_kwarg)
        by_campaigns = OrderedDict()

        # XXX `pending_for` will also include grants pending acceptance.
        requests = PortfolioDoubleOptIn.objects.pending_for(
            self.account, at_time=at_time, campaign=campaign_filtered).exclude(
                models.Q(grantee=self.account) &
                models.Q(state=PortfolioDoubleOptIn.OPTIN_GRANT_INITIATED)
        ).order_by('campaign__title')
        for optin in requests:
            campaign = optin.campaign
            if campaign:
                # XXX It is possible the request isn't limited
                #     to a single campaign.
                if not campaign in by_campaigns:
                    by_campaigns[optin.campaign] = {
                        'slug': campaign.slug,
                        'title': campaign.title,
                        'descr': campaign.description,
                        'ends_at': None,
                        'last_completed_at': None,
                        'share_url': None,
                        'respondents': [],
                        'update_url': None,
                        'requests': []}
                if optin.ends_at:
                    by_campaigns[campaign]['ends_at'] = (
                        optin.ends_at if not by_campaigns[campaign]['ends_at']
                        else min(datetime_or_now(
                            by_campaigns[campaign]['ends_at']), optin.ends_at
                    )).isoformat()
                by_campaigns[campaign]['requests'] += [{
                    'created_at': optin.created_at.isoformat(),
                    'grantee': optin.grantee.slug
                }]

        candidates = get_latest_active_assessments(
            self.account, campaign=campaign_filtered).exclude(
                campaign__slug__endswith='-verified') # XXX Ad-hoc exclude
                                                    # of verification campaigns.
        for sample in candidates:
            if not sample.campaign in by_campaigns:
                by_campaigns[sample.campaign] = {
                    'slug': sample.campaign.slug,
                    'title': sample.campaign.title,
                    'descr': sample.campaign.description,
                    'ends_at': None,
                    'last_completed_at': None,
                    'share_url': None,
                    'respondents': [],
                    'update_url': None,
                    'requests': []}
            reverse_kwargs = {
                self.account_url_kwarg: self.account,
                'sample': sample,
            }
            if path_filtered:
                reverse_kwargs.update({self.path_url_kwarg: path_filtered})
            by_campaigns[sample.campaign]['update_url'] = self.get_redirect_url(
                **reverse_kwargs)
            latest_completed = get_latest_completed_assessment(self.account,
                campaign=sample.campaign)
            if latest_completed:
                by_campaigns[sample.campaign]['last_completed_at'] = \
                    latest_completed.created_at
                by_campaigns[sample.campaign]['share_url'] = reverse(
                    'share', args=(self.account, latest_completed))
                by_campaigns[sample.campaign]['respondents'] = \
                    get_user_model().objects.filter(
                        answer__sample=sample).distinct()

        if not campaign_filtered:
            for campaign in self.campaign_candidates:
                if not campaign in by_campaigns:
                    by_campaigns[campaign] = {
                        'slug': campaign.slug,
                        'title': campaign.title,
                        'descr': campaign.description,
                        'ends_at': None,
                        'last_completed_at': None,
                        'share_url': None,
                        'respondents': [],
                        'update_url': None,
                        'requests': []}

        # We have a data structure here that looks like
        # by_campaiagns = {
        #   'Slug': {
        #     'title': String,
        #     'assess_url': URL,
        #     'last_completed_at': Datetime,
        #     'share_url': URL,
        #     'grantees': [{
        #         'slug': Slug,
        #         'created_at': Datetime
        #     }]
        #   }
        # }

        context.update({
            'candidates': by_campaigns.values()
        })
        update_context_urls(context, {
            'api_accounts': site_url("/api/profile")})
        return context


    def post(self, request, *args, **kwargs):
        form = self.get_form()
        if form.is_valid():
            return self.form_valid(form)
        return self.form_invalid(form)



class ImprovePracticesView(AssessPracticesView):
    """
    Improvement planning page
    """
    template_name = 'app/improve/index.html'
    breadcrumb_url = 'improve_practices'
    breadcrumb_url_index = 'improve_redirect'

    def get_template_names(self):
        campaign_slug = self.sample.campaign.slug
        candidates = ['app/improve/%s.html' %  campaign_slug]
        candidates += super(ImprovePracticesView, self).get_template_names()
        return candidates

    def get_context_data(self, **kwargs):
        context = super(ImprovePracticesView, self).get_context_data(**kwargs)
        context.update({'assess_type': 'planned'})
        update_context_urls(context, {
            'api_improvement_sample': reverse('survey_api_sample', args=(
                self.account, self.improvement_sample)),
            'print': reverse('improve_print', args=(
                self.account, self.improvement_sample))
        })
        if self.full_path:
            update_context_urls(context, {
                'api_account_benchmark': reverse('survey_api_sample_benchmarks',
                    args=(self.account, self.sample,
                          self.full_path.lstrip(URL_PATH_SEP))),
            })
        else:
            update_context_urls(context, {
                'api_account_benchmark': reverse(
                    'survey_api_sample_benchmarks_index', args=(
                    self.account, self.sample)),
            })
        return context


class ImproveRedirectView(ImprovePracticesView):
    """
    Index page for improvement planning
    """


class TrackMetricsView(AccountMixin, TemplateView):
    """
    Profile metrics page

    This page includes Energy tracking & GHG emissions estimator, Water use,
    and Waste tracking.

    In the Energy tracking & GHG emissions estimator case,
    the page calls `Lists profiles in a group _http://localhost:8040/djaopsp\
/api/supplier-1/filters/accounts/scope1-stationary-combustion?\
ends_at=2023-12-31&timezone=local&page=1` for group
    'scope1-stationary-combustion', 'scope1-mobile-combustion',
    'scope1-refrigerants', '712-total-scope-2-ghg-emissions',
    and '16-total-scope-3-ghg-emissions'.

    The groups are selected based on the metric tracker page URL path that
    maps to a set of question paths encoded in the `EditableFilter.extra` field.
    To differentiate path prefixes, a slug is also encoded
    in `EditableFilter.extra['tags']`.
    `EditableFilter.account` is always used in queries.


    """
    template_name = 'app/track/index.html'

    def get_template_names(self):
        candidates = []
        metric = self.kwargs.get('metric')
        if metric:
            candidates += ['app/track/%s.html' % str(metric)]
        candidates += super(TrackMetricsView, self).get_template_names()
        return candidates

    def get_editable_filter_context(self, context, candidate_paths, title=None):
        #pylint:disable=too-many-locals
        path = None
        question = None
        question_model = get_question_model()
        for candidate_path in candidate_paths:
            try:
                question = question_model.objects.get(path=candidate_path)
                path = candidate_path
                break
            except question_model.DoesNotExist:
                continue
        if not question:
            raise Http404("Cannot find one of %s" % str(candidate_paths))

        filter_args = models.Q(extra__contains=path)
        tag = None
        if not title:
            title = question.title
        else:
            tag = slugify(title)
            filter_args &= models.Q(extra__contains=tag)
        try:
            editable_filter = None
            for row in EditableFilter.objects.filter(
                filter_args, account=self.account):
                # Sometimes titles are prefixes/suffixes of each other.
                # ex: 'Hazardous Waste' and 'Non-Hazardous Waste'
                extra_path = get_extra(row, 'path', "")
                if path == extra_path:
                    if not tag:
                        editable_filter = row
                        break
                    extra_tags = get_extra(row, 'tags', [])
                    if tag in extra_tags:
                        editable_filter = row
                        break
            if not editable_filter:
                raise EditableFilter.DoesNotExist()
        except EditableFilter.DoesNotExist:
            extra = {'path': path}
            if tag:
                extra.update({'tags': [tag]})
            editable_filter = EditableFilter.objects.create(
                account=self.account,
                title=title,
                extra=json.dumps(extra))
        if tag:
            slug_part = tag
        else:
            slug_part = path.split(DB_PATH_SEP)[-1]
        api_endpoint = "api_track_%s" % slug_part.replace('-', '_')
        update_context_urls(context, {
            api_endpoint: reverse(
                'survey_api_accounts_filter',
                args=(self.account, editable_filter.slug,))
        })
        return context

    def get_context_data(self, **kwargs):
        context = super(TrackMetricsView, self).get_context_data(**kwargs)
        ends_at = datetime.date(datetime_or_now().year - 1, 12, 31)
        context.update({
            'starts_at': datetime.date(ends_at.year, 1, 1).isoformat(),
            'ends_at': ends_at.isoformat()
        })
        metric = self.kwargs.get('metric')
        #pylint:disable=line-too-long
        if metric == 'energy-ghg-emissions':
            self.get_editable_filter_context(context, ['/sustainability/environmental-reporting/ghg-emissions-reporting/ghg-emissions-scope1-measured/ghg-emissions-scope1'],
                    title='Scope1 Stationary Combustion')
            self.get_editable_filter_context(context, ['/sustainability/environmental-reporting/ghg-emissions-reporting/ghg-emissions-scope1-measured/ghg-emissions-scope1'],
                    title='Scope1 Mobile Combustion')
            self.get_editable_filter_context(context, ['/sustainability/environmental-reporting/ghg-emissions-reporting/ghg-emissions-scope1-measured/ghg-emissions-scope1'],
                    title='Scope1 Refrigerants')
            self.get_editable_filter_context(context, ['/sustainability/environmental-reporting/ghg-emissions-reporting/ghg-emissions-scope1-measured/ghg-emissions-scope2'])
            self.get_editable_filter_context(context, ['/sustainability/environmental-reporting/ghg-emissions-reporting/ghg-emissions-scope3-details/ghg-emissions-scope3-measured/ghg-emissions-scope3'])
        elif metric == 'waste':
            self.get_editable_filter_context(context, ['/sustainability/environmental-reporting/waste-reporting/waste-measured/hazardous-waste'])
            self.get_editable_filter_context(context, ['/sustainability/environmental-reporting/waste-reporting/waste-measured/non-hazardous-waste'])
            self.get_editable_filter_context(context, ['/sustainability/environmental-reporting/waste-reporting/waste-measured/waste-recycled'])
        elif metric == 'water':
            self.get_editable_filter_context(context, ['/sustainability/environmental-reporting/water-reporting/water-measured/water-withdrawn'])
            self.get_editable_filter_context(context, ['/sustainability/environmental-reporting/water-reporting/water-measured/water-discharged'])
            self.get_editable_filter_context(context, ['/sustainability/environmental-reporting/water-reporting/water-measured/water-recycled'])
        return context


class AssessPracticesXLSXView(AssessmentContentMixin, PracticesSpreadsheetView):

    base_headers = ['', 'Assessed', 'Planned', 'Comments', 'Opportunity']

    @property
    def intrinsic_value_headers(self):
        if not hasattr(self, '_intrinsic_value_headers'):
            #pylint:disable=attribute-defined-outside-init
            self._intrinsic_value_headers = []
            for seg in self.segments_available:
                prefix = seg.get('path')
                if prefix:
                    score_calculator = get_score_calculator(prefix)
                    if (score_calculator and
                        score_calculator.intrinsic_value_headers):
                        self._intrinsic_value_headers = \
                            score_calculator.intrinsic_value_headers
                        break
        return self._intrinsic_value_headers

    def get_headings(self):
        if not hasattr(self, 'units'):
            #pylint:disable=attribute-defined-outside-init
            self.units = {}
        self.peer_value_headers = []
        for unit in six.itervalues(self.units):
            if (unit.system == unit.SYSTEM_ENUMERATED and
                unit.slug == 'assessment'):
                self.peer_value_headers += [
                    (unit.slug, [choice.text for choice in unit.choices])]
        headers = [] + self.base_headers
        for header in self.peer_value_headers:
            headers += header[1]
        if self.peer_value_headers:
            headers += [
                "Nb respondents",
            ]
        headers += self.intrinsic_value_headers
        return headers

    def get_title_row(self):
        return [self.sample.account.printable_name,
            self.sample.created_at.strftime("%Y-%m-%d")]


    def format_row(self, entry, key=None):
        #pylint:disable=too-many-locals
        default_unit = entry.get('default_unit', {})
        default_unit_choices = []
        if default_unit:
            try:
                default_unit = default_unit.slug
            except AttributeError:
                default_unit_choices = default_unit.get('choices', [])
                default_unit = default_unit.get('slug', "")
        answers = entry.get('answers')
        primary_assessed = None
        primary_planned = None
        points = None
        comments = ""
        if answers:
            for answer in answers:
                unit = answer.get('unit')
                if unit and default_unit and unit.slug == default_unit:
                    primary_assessed = answer.get('measured')
                    continue
                if unit and unit.slug == 'freetext': #XXX
                    comments = answer.get('measured')
                if unit and unit.slug == 'points': #XXX
                    points = float(answer.get('measured'))
        planned = entry.get('planned')
        if planned:
            for answer in planned:
                unit = answer.get('unit')
                if unit and default_unit and unit.slug == default_unit:
                    primary_planned = answer.get('measured')
        # base_headers
        title = entry['title']
        if default_unit and default_unit_choices:
            subtitle = ""
            for choice in default_unit_choices:
                text = choice.get('text', "").strip()
                descr = choice.get('descr', "").strip()
                if text != descr:
                    subtitle += "\n%s - %s" % (text, descr)
            if subtitle:
                title += "\n" + subtitle
        row = [
            title,
            primary_assessed,
            primary_planned,
            comments,
            entry.get('opportunity')]

        # peer_value_headers
        for header in self.peer_value_headers:
            if default_unit and default_unit == header[0]:
                for text in header[1]:
                    row += [entry.get('rate', {}).get(text, 0)]
            else:
                for text in header[1]:
                    row += [""]
        if self.peer_value_headers:
            row += [
                entry.get('nb_respondents'),
            ]

        # intrinsic_value_headers
        avg_value = 0
        extra = entry.get('extra')
        if extra:
            intrinsic_values = extra.get('intrinsic_values')
            if intrinsic_values:
                environmental_value = intrinsic_values.get('environmental', 0)
                business_value = intrinsic_values.get('business', 0)
                profitability = intrinsic_values.get('profitability', 0)
                implementation_ease = intrinsic_values.get(
                    'implementation_ease', 0)
                avg_value = (environmental_value + business_value +
                    profitability + implementation_ease) // 4
        if avg_value:
            row += [
                environmental_value,
                business_value,
                profitability,
                implementation_ease,
                avg_value
            ]
        else:
            row += ['', '', '', '', '']
        return row

    def get_filename(self):
        return datetime_or_now().strftime("%s-%s-%%Y%%m%%d.xlsx" % (
            self.sample.account.slug, self.campaign.slug))


class AssessPracticesPPTXView(AssessmentContentMixin, ListView):
    """
    Allows downloading an assessment as a PPTX
    """
    #pylint:disable=too-many-instance-attributes

    def __init__(self):
        super().__init__()
        self.prs = Presentation()
        self.template_config = {
            "background_color": [255, 255, 255],
            "title_font": "Arial",
            "title_font_size": 18,
            "title_font_color": [42, 87, 141],
            "body_font": "Calibri",
            "body_font_size": 12,
            "left": 1,
            "top": 2,
            "width": 8,
            "height": 0.5,
        }
        self.update_template_settings()
        self.current_slide = None
        self.title_hierarchy = {
            0: None, 1: None, 2: None, 3: None, 4: None, 5: None}
        self.basename = 'practices'

    def update_template_settings(self):
        self.left = Inches(self.template_config["left"])
        self.top = Inches(self.template_config["top"])
        self.width = Inches(self.template_config["width"])
        self.height = Inches(self.template_config["height"])

    def apply_background(self, slide):
        bg_color = self.template_config.get("background_color")
        if bg_color:
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = RGBColor(*bg_color)

    def format_slide_title(self, title_shape):
        title_shape.text_frame.paragraphs[0].font.name = \
            self.template_config["title_font"]
        title_shape.text_frame.paragraphs[0].font.size = Pt(
            self.template_config["title_font_size"])
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(
            *self.template_config["title_font_color"])

    def format_entry_content(self, entry):
        #pylint:disable=too-many-locals
        content = []
        try:
            default_unit_slug = entry['default_unit'].slug
        except AttributeError:
            default_unit_slug = entry.get('default_unit', {}).get('slug', "")

        answers = entry.get('answers', [])
        primary_assessed = None
        primary_planned = None
        points = None
        comments = ""

        for answer in answers:
            try:
                unit_slug = answer['unit'].slug
            except AttributeError:
                unit_slug = ""

            measured = answer.get('measured')
            if unit_slug == default_unit_slug:
                primary_assessed = measured
            elif unit_slug == 'freetext':
                comments = measured
            elif unit_slug == 'points':
                points = float(measured) if measured else None

        planned = entry.get('planned', [])
        for plan in planned:
            try:
                unit_slug = plan['unit'].slug
            except AttributeError:
                unit_slug = ""

            if unit_slug == default_unit_slug:
                primary_planned = plan.get('measured')

        title = entry['title']
        content.append(f"Title: {title}")
        opportunity = entry.get('opportunity')
        rates = entry.get('rate', {})
        nb_respondents = entry.get('nb_respondents', {})
        extra = entry.get('extra')
        avg_value = 0

        if primary_assessed is not None:
            content.append(f"Assessed: {primary_assessed}")
        if primary_planned is not None:
            content.append(f"Planned: {primary_planned}")
        if points is not None:
            content.append(f"Points: {points}")
        if comments:
            content.append(f"Comments: {comments}")
        if opportunity:
            content.append(f"Opportunity: {opportunity}")
        if nb_respondents:
            content.append(f"Number of respondents: {nb_respondents}")
        if rates:
            for choice, value in rates.items():
                content.append(f"{choice}: {value}")
        if extra:
            intrinsic_values = extra.get('intrinsic_values')
            if intrinsic_values:
                environmental_value = intrinsic_values.get('environmental', 0)
                business_value = intrinsic_values.get('business', 0)
                profitability = intrinsic_values.get('profitability', 0)
                implementation_ease = intrinsic_values.get(
                    'implementation_ease', 0)
                avg_value = (environmental_value + business_value +
                    profitability + implementation_ease) // 4
            if avg_value:
                content.append(f"Environmental value: {environmental_value}")
                content.append(f"Business value: {business_value}")
                content.append(f"Profitability: {profitability}")
                content.append(f"Implementation Ease: {implementation_ease}")
                content.append(f"Average Value: {avg_value}")

        return "\n".join(content)


    def add_new_slide(self, presentation, title):
        #pylint:disable=attribute-defined-outside-init
        self.top = Inches(self.template_config["top"])
        slide_layout = presentation.slide_layouts[5]
        slide = presentation.slides.add_slide(slide_layout)
        self.apply_background(slide)
        title_shape = slide.shapes.title
        title_shape.text = title
        self.format_slide_title(title_shape)
        return slide

    def add_content_to_slide(self, slide, entry):
        textbox = slide.shapes.add_textbox(self.left, self.top, self.width, self.height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        content = self.format_entry_content(entry)

        for line in content.split("\n"):
            para = text_frame.add_paragraph()

            if ": " in line:
                prefix, rest = line.split(": ", 1)
                self.add_text_run(para, prefix + ":", bold=True)
            else:
                rest = line

            self.add_text_run(para, rest)

    def add_text_run(self, paragraph, text, bold=False):
        run = paragraph.add_run()
        run.text = text
        run.font.bold = bold
        run.font.name = self.template_config["body_font"]
        run.font.size = Pt(self.template_config["body_font_size"])

    def add_extra_content_to_title_slide(self, slide, extra_values):
        for key, value in extra_values.items():
            textbox = slide.shapes.add_textbox(
                self.left, self.top, self.width, self.height)
            text_frame = textbox.text_frame
            text_frame.word_wrap = True

            if key and value:
                para = text_frame.add_paragraph()
                para.text = f"{key}: {value}"
                for run in para.runs:
                    run.font.name = self.template_config["body_font"]
                    run.font.size = Pt(self.template_config["body_font_size"])

            self.top += self.height

    def get(self, request, *args, **kwargs):
        self.current_slide = None
        queryset = self.get_queryset()

        for entry in queryset:
            indent_level = entry['indent']
            title = entry['title']

            if indent_level > 1:
                self.title_hierarchy[indent_level] = title
                for higher_level in range(indent_level + 1, 6):
                    self.title_hierarchy[higher_level] = None

                if 'answers' in entry:
                    title_parts = [self.title_hierarchy[level]
                        for level in range(1, indent_level)
                        if self.title_hierarchy[level] is not None]
                    composed_title = ' - '.join(title_parts)
                    self.current_slide = self.add_new_slide(
                        self.prs, composed_title)
                    self.add_content_to_slide(self.current_slide, entry)

            elif indent_level == 0:
                self.current_slide = self.add_new_slide(self.prs, title)
                extra_values = {
                    'Campaign': self.sample.campaign.title,
                    'Created at': self.sample.created_at.date(),
                    'Campaign Creator': self.sample.account.full_name,
                    'Normalized score': entry.get('normalized_score'),
                }
                self.add_extra_content_to_title_slide(
                    self.current_slide, extra_values)

            elif indent_level == 1:
                self.current_slide = self.add_new_slide(self.prs, title)
                self.title_hierarchy[1] = title
                for higher_level in range(2, 6):
                    self.title_hierarchy[higher_level] = None
        ppt_io = BytesIO()
        self.prs.save(ppt_io)
        ppt_io.seek(0)
        filename = self.get_filename()
        response = HttpResponse(
            ppt_io,
            content_type=self.content_type)
        response['Content-Disposition'] = f'attachment; filename={filename}'
        return response

    def get_filename(self):
        return datetime_or_now().strftime("%s-%s-%%Y%%m%%d.pptx" % (
            self.sample.account.slug, self.campaign.slug))
